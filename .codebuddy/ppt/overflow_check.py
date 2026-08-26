# -*- coding: utf-8 -*-
# 文本溢出估算校验：估算每个文本框内容所需高度/宽度，与实际框大小对比
import sys, io, re
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu

EMU = 914400.0

def char_w(ch, fs):
    # 中日韩全角字符 ≈ fs/72 英寸；半角 ≈ 0.52*fs/72
    if re.match(r'[\u2e80-\u9fff\uff00-\uffef\u3000-\u303f\u2018\u2019\u201c\u201d\u2014\u2026\u2192\u2190]', ch):
        return fs / 72.0
    return 0.52 * fs / 72.0

def text_width(s, fs):
    return sum(char_w(c, fs) for c in s)

prs = Presentation("热重载方案分享.pptx")
warns = []
for idx, slide in enumerate(prs.slides, 1):
    for shp in slide.shapes:
        if not shp.has_text_frame:
            continue
        tf = shp.text_frame
        if not tf.text.strip():
            continue
        box_w = (shp.width or 0) / EMU
        box_h = (shp.height or 0) / EMU
        if box_w <= 0 or box_h <= 0:
            continue
        need_h = 0.0
        max_overrun_line = None
        for para in tf.paragraphs:
            # 取该段字号（run 级，取最大）
            fs = 12.0
            for r in para.runs:
                if r.font.size:
                    fs = max(fs, r.font.size.pt)
            ptxt = "".join(r.text for r in para.runs)
            if not ptxt:
                need_h += fs / 72.0 * 1.2
                continue
            # bullet 缩进
            indent = 0.18 if para._pPr is not None and para._pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buChar') is not None else 0
            avail = box_w - indent - 0.05
            w = text_width(ptxt, fs)
            lines = max(1, int(w / avail) + (1 if w % avail > 0.01 else 0))
            line_h = fs / 72.0 * 1.30
            # 段后距
            spa = 0.0
            try:
                if para.space_after:
                    spa = para.space_after.pt / 72.0
            except Exception:
                pass
            need_h += lines * line_h + spa
            if lines > 1 and max_overrun_line is None:
                max_overrun_line = ptxt[:30]
        if need_h > box_h * 1.06:  # 6% 容差
            warns.append(f"页{idx}: 文本可能溢出 框高{box_h:.2f}\" 需≈{need_h:.2f}\" | 首个换行段: {max_overrun_line or tf.text[:30]!r}")

print("===== 文本溢出估算 =====")
print("\n".join(warns) if warns else "未发现明显溢出")
