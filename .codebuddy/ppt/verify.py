# -*- coding: utf-8 -*-
# 校验 PPT：页数、每页文本、检查占位残留与关键内容
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu

prs = Presentation("热重载方案分享.pptx")
print(f"总页数: {len(prs.slides)}")
print(f"画布: {prs.slide_width/914400:.2f} x {prs.slide_height/914400:.2f} in")
issues = []
for idx, slide in enumerate(prs.slides, 1):
    texts = []
    for shp in slide.shapes:
        if shp.has_text_frame:
            t = shp.text_frame.text.strip().replace("\n", " / ")
            if t:
                texts.append(t)
        # 越界检查
        try:
            if shp.left is not None and shp.top is not None:
                r = (shp.left + (shp.width or 0)) / 914400
                b = (shp.top + (shp.height or 0)) / 914400
                if r > 13.34 or b > 7.51 or shp.left < -1000 or shp.top < -1000:
                    issues.append(f"页{idx}: 形状越界 right={r:.2f} bottom={b:.2f}")
        except Exception:
            pass
    joined = " || ".join(texts)
    print(f"\n--- 页 {idx} ---")
    print(joined[:400])
    for bad in ["lorem", "ipsum", "xxxx", "TODO", "undefined"]:
        if bad.lower() in joined.lower():
            issues.append(f"页{idx}: 残留占位文本 {bad}")

print("\n\n===== 问题 =====")
print("\n".join(issues) if issues else "无")
